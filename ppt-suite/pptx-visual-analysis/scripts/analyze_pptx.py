#!/usr/bin/env python3
"""通用 PPTX 视觉探针：打印画布尺寸、主题色 scheme、字体/字号/颜色 Counter。

用法: /usr/bin/python3 analyze_pptx.py <file.pptx>
注意: 用 /usr/bin/python3（自带 pptx + PIL + pymupdf），shell 默认 venv python 无 pptx。
"""
import sys, re
from collections import Counter
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


def color_str(color):
    """鲁棒读取 run 颜色；scheme 色无 .rgb 会抛 AttributeError。"""
    try:
        t = color.type
        if t is None:
            return 'none'
        tn = str(t)
        if 'SCHEME' in tn.upper():
            return f'scheme:{color.theme_color}'
        if 'RGB' in tn.upper():
            return str(color.rgb)
        return tn
    except Exception:
        return f'?{type(color).__name__}'


def theme_colors(prs):
    """主题色从 theme part + regex 取（element.find('.//a:clrScheme') 会返回 None）。"""
    for mi, master in enumerate(prs.slide_masters):
        try:
            tp = master.part.part_related_by(
                'http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme')
            xml = tp.blob.decode('utf-8', errors='ignore')
            m = re.search(r'<a:clrScheme.*?</a:clrScheme>', xml, re.S)
            if not m:
                continue
            print(f'master[{mi}] clrScheme:')
            for tag in ['dk1', 'lt1', 'dk2', 'lt2',
                        'accent1', 'accent2', 'accent3', 'accent4', 'accent5', 'accent6']:
                mm = re.search(rf'<a:{tag}>\s*<a:srgbClr val="([0-9A-Fa-f]{{6}})"', m.group(0))
                if mm:
                    print(f'  {tag}: #{mm.group(1)}')
        except Exception as e:
            print(f'master[{mi}] theme err: {e}')


def analyze(path):
    prs = Presentation(path)
    print(f'SLIDE SIZE: {Emu(prs.slide_width).inches:.2f}" x {Emu(prs.slide_height).inches:.2f}"'
          f'  ({Emu(prs.slide_width).cm:.1f} x {Emu(prs.slide_height).cm:.1f} cm)')
    print(f'TOTAL SLIDES: {len(prs.slides)}')

    theme_colors(prs)

    fonts = Counter()
    sizes = Counter()
    colors = Counter()

    def walk(sh):
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts[run.font.name] += 1
                    if run.font.size:
                        sizes[run.font.size.pt] += 1
                    if run.font.color and run.font.color.type is not None:
                        colors[color_str(run.font.color)] += 1
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in sh.shapes:
                walk(sub)

    for slide in prs.slides:
        for sh in slide.shapes:
            walk(sh)

    print('\nFONTS:', dict(fonts.most_common(12)))
    print('SIZES:', dict(sizes.most_common(12)))
    print('COLORS:', dict(colors.most_common(12)))
    return prs


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    analyze(sys.argv[1])
